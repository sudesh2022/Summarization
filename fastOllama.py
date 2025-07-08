# app.py

import os
import re
import asyncio
import hashlib
from typing import List
from fastapi import FastAPI, File, UploadFile, Request
from fastapi.responses import HTMLResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from starlette.middleware.cors import CORSMiddleware
from pptx import Presentation
from pdfminer.high_level import extract_text
from openai import AsyncOpenAI
from dotenv import load_dotenv
from databases import Database
import sqlalchemy
from sqlalchemy import create_engine
from contextlib import asynccontextmanager
import ollama

# === Load API Key ===
#load_dotenv()
#client = AsyncOpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# === FastAPI Setup ===
app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.mount("/static", StaticFiles(directory="static"), name="static")
templates = Jinja2Templates(directory="templates")

# === Database Setup ===

DATABASE_URL = "postgresql+asyncpg://sudesh:sudesh@localhost:5432/summariesdb"
database = Database(DATABASE_URL)
metadata = sqlalchemy.MetaData()

summaries_table = sqlalchemy.Table(
    "summaries",
    metadata,
    sqlalchemy.Column("id", sqlalchemy.Integer, primary_key=True),
    sqlalchemy.Column("filename", sqlalchemy.String),
    sqlalchemy.Column("filehash", sqlalchemy.String, unique=True),
    sqlalchemy.Column("summary", sqlalchemy.Text),
)

# Sync engine for initial table creation (once)
sync_engine = create_engine(DATABASE_URL.replace("+asyncpg", ""))
metadata.create_all(sync_engine)

@asynccontextmanager
async def lifespan(app):
    await database.connect()
    yield
    await database.disconnect()

app = FastAPI(lifespan=lifespan)

# === DB Logic ===
async def get_summary_from_db(file_hash: str):
    query = summaries_table.select().where(summaries_table.c.filehash == file_hash)
    row = await database.fetch_one(query)
    return row["summary"] if row else None

async def save_summary_to_db(filename: str, file_hash: str, summary: str):
    query = summaries_table.insert().values(
        filename=filename,
        filehash=file_hash,
        summary=summary
    )
    await database.execute(query)

# === Utility Functions ===

def file_hash(file_path: str) -> str:
    hasher = hashlib.sha256()
    with open(file_path, "rb") as f:
        hasher.update(f.read())
    return hasher.hexdigest()

def smart_chunk(text: str, max_words=300) -> list:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    chunks, current = [], []
    for line in lines:
        current.append(line)
        if sum(len(w.split()) for w in current) >= max_words:
            chunks.append(" ".join(current))
            current = []
    if current:
        chunks.append(" ".join(current))
    return chunks

def read_pdf(file_path: str) -> list:
    text = extract_text(file_path)
    return smart_chunk(text)

def read_ppt(file_path: str) -> list:
    prs = Presentation(file_path)
    slides = "\n".join([
        shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
    ])
    return smart_chunk(slides)

async def query_chatgpt4(prompt: str) -> str:
    try:
        response = await client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5,
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"Error querying ChatGPT-4: {e}")
        return ""

async def query_ollama(prompt: str, model: str = "deepseek-r1") -> str:
    try:
        response = await asyncio.to_thread(
            ollama.chat,
            model=model,
            messages=[{"role": "user", "content": prompt}]
        )
        return response['message']['content'].strip()
    except Exception as e:
        print(f"Error querying Ollama: {e}")
        return ""

async def summarize_chunk(chunk: str, model: str = "deepseek-r1") -> str:
    prompt = (
        "You are a summarization system. Generate a short summary of the following text. "
        "Clarify the main points in a clear structure. Capture key ideas and essential information. "
        "Strictly follow grammatical and syntactical rules.\n\n"
        f"{chunk}"
    )
    return await query_ollama(prompt, model)

async def map_reduce_summary(chunks: list, model: str = "deepseek-r1") -> str:
    summaries = await asyncio.gather(*(summarize_chunk(chunk, model) for chunk in chunks))
    return await summarize_chunk(" ".join(summaries), model)

def format_summary(summary: str) -> str:
    lines = summary.splitlines()
    formatted = []
    bullet_points = []
    for line in lines:
        match = re.match(r"^(\d+)\.\s+(.*)", line)
        if match:
            bullet_points.append(f"<li>{match.group(2)}</li>")
        else:
            if bullet_points:
                formatted.append("<ul>" + "".join(bullet_points) + "</ul>")
                bullet_points = []
            if line.strip():
                formatted.append(f"<p>{line.strip()}</p>")
    if bullet_points:
        formatted.append("<ul>" + "".join(bullet_points) + "</ul>")
    return "".join(formatted)

# === Route ===
@app.api_route("/", methods=["GET", "POST"], response_class=HTMLResponse)
async def index(request: Request, files: List[UploadFile] = File(None)):
    summaries = []

    if request.method == "POST" and files:
        for file in files:
            file_path = os.path.join(UPLOAD_FOLDER, file.filename)
            with open(file_path, "wb") as f:
                f.write(await file.read())

            f_hash = file_hash(file_path)

            # 1. Check DB for cached summary
            cached_summary = await get_summary_from_db(f_hash)
            if cached_summary:
                summaries.append({"filename": file.filename, "summary": cached_summary})
                continue

            # 2. Process new file
            ext = file.filename.lower()
            if ext.endswith(".pdf"):
                text_chunks = await asyncio.to_thread(read_pdf, file_path)
            elif ext.endswith(".ppt") or ext.endswith(".pptx"):
                text_chunks = await asyncio.to_thread(read_ppt, file_path)
            else:
                summaries.append({"filename": file.filename, "summary": "Unsupported file type."})
                continue

            if text_chunks:
                summary_text = await map_reduce_summary(text_chunks, model="deepseek-r1")
                formatted = format_summary(summary_text)

                # Save to DB
                await save_summary_to_db(file.filename, f_hash, formatted)

                summaries.append({"filename": file.filename, "summary": formatted})
            else:
                summaries.append({"filename": file.filename, "summary": "No text found in file."})

    return templates.TemplateResponse("fastindex.html", {"request": request, "summaries": summaries})
