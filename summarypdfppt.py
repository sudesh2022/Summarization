import ollama
from pdfminer.high_level import extract_text
from pptx import Presentation
import concurrent.futures
import time
import os
import re
from flask import Flask, render_template, request

app = Flask(__name__)
UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

def query_ollama(model: str, prompt: str):
    """Query Ollama API with a given model and prompt."""
    try:
        response = ollama.chat(
            model=model,
            messages=[{"role": "user", "content": prompt}]
        )
        return response.get("message", {}).get("content", "No response")
    except Exception as e:
        print(f"Error querying Ollama: {e}")
        return ""

def read_pdf(file_path: str) -> list:
    """Extract text from a PDF efficiently using pdfminer."""
    text = extract_text(file_path)
    text_chunks = text.split("\n\n")  # Split into chunks by paragraphs
    print(f"Extracted {len(text_chunks)} text chunks from {file_path}.")
    return text_chunks

def read_ppt(file_path: str) -> list:
    """Extract text from a PowerPoint file."""
    prs = Presentation(file_path)
    text_chunks = []
    for slide in prs.slides:
        slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        text_chunks.append(slide_text)
    print(f"Extracted {len(text_chunks)} text chunks from {file_path}.")
    return text_chunks

def summarize_chunk(chunk: str, model: str) -> str:
    """Summarize a chunk of text using Ollama."""
    print("Summarizing chunk...")
    mySummaryPrompt1 = ''' You are a english summarization system. Generate a general summary of the following text in 2000 words. 
    The result must strictly adhere to the grammatical and syntactical rules of language. 
    Any deviation from the specified language will not be acceptable. 
    Clarify the main topic of text in a clear structure. Do not provide explanations.
    Please note, you should not invent any information. 
    Stick to the facts provided in the texts.
    '''
    mySummaryPrompt = '''You are a summarization system. Generate a short summary in  of the following text. Clarify the main points of text in a clear structure. If text is provided, capture the main points, key ideas, and essential information in text. . MUST strictly adhere to the grammatical and syntactical rules. Any deviation from the specified language will not be acceptable'''

    prompt = f"{mySummaryPrompt}:\n{chunk}"
    return query_ollama(model, prompt)

def map_reduce_summary(chunks: list, model: str) -> str:
    """Perform map-reduce summarization using parallel execution."""
    summaries = []
    
    # Use more threads for parallel execution
    with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
        future_to_chunk = {executor.submit(summarize_chunk, chunk, model): chunk for chunk in chunks}
        for future in concurrent.futures.as_completed(future_to_chunk):
            try:
                summaries.append(future.result())
            except Exception as e:
                print(f"Error processing chunk: {e}")
    
    print("Combining partial summaries...")
    combined_text = " ".join(summaries)  # Use space instead of newline for efficiency
    return summarize_chunk(combined_text, model)  # Reduce step

def format_summary(summary: str) -> list:
    """Format summary into a structured table with numbered points."""
    lines = summary.split("\n")
    formatted_summary = []
    for line in lines:
        match = re.match(r"^(\d+)\.\s(.+)", line)
        if match:
            formatted_summary.append({"point": match.group(1), "text": match.group(2)})
        else:
            formatted_summary.append({"point": "", "text": line})
    return formatted_summary

@app.route("/", methods=["GET", "POST"])
def upload_files():
    summaries = []
    if request.method == "POST":
        files = request.files.getlist("files")
        for file in files:
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(file_path)
            
            text_chunks = []
            if file.filename.endswith(".pdf"):
                with concurrent.futures.ProcessPoolExecutor() as executor:
                    future = executor.submit(read_pdf, file_path)
                    text_chunks = future.result()
            elif file.filename.endswith(".ppt") or file.filename.endswith(".pptx"):
                with concurrent.futures.ProcessPoolExecutor() as executor:
                    future = executor.submit(read_ppt, file_path)
                    text_chunks = future.result()
            
            if text_chunks:
                summary_text = map_reduce_summary(text_chunks, "mistral")
                formatted_summary = format_summary(summary_text)
                summaries.append({"filename": file.filename, "summary": formatted_summary})
            else:
                summaries.append({"filename": file.filename, "summary": "No text found in file."})
    
    return render_template("index1.html", summaries=summaries)

if __name__ == "__main__":
    app.run(debug=True)